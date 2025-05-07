from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.fill import FillFormat
from pptx.enum.dml import MSO_FILL

def create_garbotgpt_presentation():
    # Initialize presentation
    prs = Presentation()
    prs.slide_width = Inches(16)  # Widescreen 16:9
    prs.slide_height = Inches(9)

    # Color palette
    bg_color = RGBColor(10, 15, 28)  # Dark blue (#0A0F1C)
    text_color = RGBColor(232, 235, 252)  # Light gray (#E8EBFC)
    accent_color = RGBColor(0, 228, 255)  # Cyan (#00E4FF)
    secondary_color = RGBColor(123, 44, 191)  # Purple (#7B2CBF)

    # Helper function to set gradient background
    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_stops[0].color.rgb = bg_color
        fill.gradient_stops[1].color.rgb = secondary_color
        fill.gradient_angle = 45

    # Helper function to add text box
    def add_text_box(slide, text, left, top, width, height, font_size, color, bold=False, align=PP_ALIGN.LEFT):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = text
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.runs[0] if p.runs else p.add_run()
        run.text = text
        font = run.font
        font.name = 'Calibri'
        font.size = Pt(font_size)
        font.color.rgb = color
        font.bold = bold
        return txBox

    # Helper function to add code box
    def add_code_box(slide, code, left, top, width, height, font_size):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(30, 30, 30)  # Dark code background
        shape.line.color.rgb = text_color
        txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
        tf = txBox.text_frame
        tf.text = code
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.runs[0] if p.runs else p.add_run()
        run.text = code
        font = run.font
        font.name = 'Consolas'
        font.size = Pt(font_size)
        font.color.rgb = RGBColor(200, 200, 200)
        return shape

    # Helper function to add image with border and title
    def add_image_with_border(slide, img_path, left, top, width, height, img_title):
        # Add image title
        add_text_box(slide, img_title, left, top - Inches(0.5), width, Inches(0.5), 16, accent_color, bold=True, align=PP_ALIGN.CENTER)
        # Add a rectangular border
        border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left - Inches(0.1), top - Inches(0.1),
            width + Inches(0.2), height + Inches(0.2)
        )
        border.fill.solid()
        border.fill.fore_color.rgb = accent_color
        border.line.color.rgb = text_color
        # Add the image
        try:
            slide.shapes.add_picture(img_path, left, top, width, height)
        except FileNotFoundError:
            add_text_box(slide, f"[Image Missing: {img_path}]", left, top, width, height, 16, text_color)

    # Helper function to add icon
    def add_icon(slide, left, top, size, color):
        shape = slide.shapes.add_shape(
            MSO_SHAPE.HEXAGON, left, top, size, size
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = text_color
        return shape

    # Slide 1: Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    title = slide.shapes.title
    title.text = "GarBotGPT: Creació d’un Chatbot d’IA"
    title.text_frame.paragraphs[0].font.size = Pt(60)
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True

    subtitle = slide.placeholders[1]
    subtitle.text = "Un projecte innovador de GarolaCorp"
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle.text_frame.paragraphs[0].font.color.rgb = text_color

    try:
        slide.shapes.add_picture(
            'static/img/garbotgpt-logo.png',
            left=Inches(12.5), top=Inches(0.5), height=Inches(2)
        )
    except FileNotFoundError:
        add_text_box(slide, "[Logo Placeholder]", Inches(12.5), Inches(0.5), Inches(3), Inches(0.5), 20, text_color)
    add_text_box(slide, "Nota: Afegir animació 'Zoom' al títol", Inches(0.5), Inches(8), Inches(5), Inches(0.5), 12, text_color)

    # Slide 2: Introduction (Updated with Explanations)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    title = slide.shapes.title
    title.text = "Introducció al Projecte"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True

    content = slide.shapes.placeholders[1].text_frame
    content.text = (
        "GarBotGPT és un chatbot d’IA avançat:\n\n"
        "• Objectiu: Facilitar aprenentatge i productivitat.\n"
        "  Ajuda els usuaris a resoldre dubtes ràpidament.\n"
        "• Tecnologia: Flask, React, PostgreSQL, OpenAI API.\n"
        "  Combina backend robust amb interfície dinàmica.\n"
        "• Funcionalitats: Suport multilingüe, entrada per veu.\n"
        "  Personalització i accessibilitat per a tothom.\n"
        "• Desplegament: Plataforma escalable a Render.\n"
        "  Assegura disponibilitat i rendiment al núvol."
    )
    for p in content.paragraphs:
        p.font.size = Pt(22)
        p.font.color.rgb = text_color
    add_text_box(slide, "Nota: Animació 'Fade In' per al contingut", Inches(0.5), Inches(8), Inches(5), Inches(0.5), 12, text_color)

    # Slide 3: Technologies Used (Updated with Explanations)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    title = slide.shapes.title
    title.text = "Tecnologies Utilitzades"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True

    techs = [
        ("Python", "Flask 2.0.1, openai 1.68.2", Inches(1.5)),
        ("Frontend", "Bootstrap 5, Animate.css", Inches(3)),
        ("Docker", "Contenidors amb python:3.9-slim", Inches(4.5)),
        ("PostgreSQL", "Base de dades amb psycopg2", Inches(6))
    ]

    for i, (name, desc, top) in enumerate(techs):
        add_icon(slide, Inches(1), top, Inches(0.6), secondary_color)
        add_text_box(slide, name, Inches(2), top, Inches(4), Inches(0.5), 26, accent_color, bold=True)
        explanation = {
            "Python": "Backend amb Flask i integració amb OpenAI.",
            "Frontend": "Interfície responsiva amb estils moderns.",
            "Docker": "Entorns consistents per a desenvolupament.",
            "PostgreSQL": "Gestió de dades amb alta fiabilitat."
        }[name]
        add_text_box(slide, f"{desc}\n  {explanation}", Inches(2), top + Inches(0.5), Inches(8), Inches(0.7), 20, text_color)
    add_text_box(slide, "Nota: Animació 'Slide In' per a tecnologies", Inches(0.5), Inches(8), Inches(5), Inches(0.5), 12, text_color)

    # Slide 4: System Architecture (Updated with Explanations)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    title = slide.shapes.title
    title.text = "Arquitectura del Sistema"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True

    content = slide.shapes.placeholders[1].text_frame
    content.text = (
        "Estructura modular:\n\n"
        "• Frontend: Bootstrap 5 i JavaScript per a la UI.\n"
        "  Interfície dinàmica i accessible per a usuaris.\n"
        "• Backend: Flask amb API RESTful.\n"
        "  Gestiona lògica i comunicació amb l’IA.\n"
        "• Base de Dades: PostgreSQL per a usuaris i xats.\n"
        "  Emmagatzema dades de manera estructurada.\n"
        "• Contenidors: Docker per a consistència.\n"
        "  Simplifica el desplegament i proves."
    )
    for p in content.paragraphs:
        p.font.size = Pt(22)
        p.font.color.rgb = text_color
    add_text_box(slide, "[Afegir diagrama: Frontend -> Backend -> DB]", Inches(10), Inches(2), Inches(5), Inches(0.5), 20, text_color)
    add_text_box(slide, "Nota: Animació 'Fade In' per al contingut", Inches(0.5), Inches(8), Inches(5), Inches(0.5), 12, text_color)

    # Slide 5: Docker Configuration (Updated with Explanations)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    title = slide.shapes.title
    title.text = "Configuració de Docker"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True

    content = slide.shapes.placeholders[1].text_frame
    content.text = (
        "• Dockerfile: Entorn amb python:3.9-slim.\n"
        "  Defineix l'entorn base per a l'aplicació.\n"
        "• docker-compose.yml: Serveis app i db.\n"
        "  Orquestra contenidors per a backend i DB.\n"
        "• PostgreSQL: Healthcheck per fiabilitat.\n"
        "  Garanteix que la DB estigui operativa.\n"
        "• Avantatges: Portabilitat entre entorns.\n"
        "  Facilita proves i desplegament al núvol."
    )
    for p in content.paragraphs:
        p.font.size = Pt(22)
        p.font.color.rgb = text_color
    code = (
        "FROM python:3.9-slim\n"
        "WORKDIR /app\n"
        "COPY requirements.txt .\n"
        "RUN pip install -r requirements.txt\n"
        "COPY . .\n"
        "CMD [\"gunicorn\", \"--bind\", \"0.0.0.0\", \"app:app\"]"
    )
    add_code_box(slide, code, Inches(9), Inches(2), Inches(6), Inches(4), 14)
    add_text_box(slide, "Nota: Animació 'Slide In' per al codi", Inches(0.5), Inches(8), Inches(5), Inches(0.5), 12, text_color)

    # Slide 6: Deployment to Render (Already Updated)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    title = slide.shapes.title
    title.text = "Desplegament a Render"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True

    content = slide.shapes.placeholders[1].text_frame
    content.text = (
        "• Repositori: GitHub amb Procfile.\n"
        "  Defineix com iniciar l'aplicació a Render.\n"
        "• Integració: Desplegament automàtic.\n"
        "  Actualitza l'aplicació amb cada canvi al codi.\n"
        "• Configuració: Variables d’entorn (.env).\n"
        "  Gestiona claus secretes com API keys.\n"
        "• Resultat: garbotgpt.com accessible.\n"
        "  Domini personalitzat amb certificat SSL."
    )
    for p in content.paragraphs:
        p.font.size = Pt(22)
        p.font.color.rgb = text_color
    code = "web: gunicorn app:app"
    add_code_box(slide, code, Inches(10), Inches(2), Inches(5), Inches(1), 18)
    add_image_with_border(slide, 'screenshots/render_dashboard_app.png', Inches(10), Inches(3.5), Inches(5), Inches(3), "Tauler de Control de l'Aplicació a Render")
    add_text_box(slide, "Panell de control de l'aplicació a Render", Inches(10), Inches(6.5), Inches(5), Inches(0.5), 16, text_color)
    add_text_box(slide, "Nota: Animació 'Fade In' per al contingut", Inches(0.5), Inches(8), Inches(5), Inches(0.5), 12, text_color)

    # Slide 7: Database Setup (Updated with Explanations)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    title = slide.shapes.title
    title.text = "Base de Dades"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True

    content = slide.shapes.placeholders[1].text_frame
    content.text = (
        "• Tecnologia: PostgreSQL 13.\n"
        "  Sistema robust per a gestió de dades.\n"
        "• Taules: User, ChatSettings, ChatHistory.\n"
        "  Emmagatzema usuaris i converses.\n"
        "• Connexió: Via docker-compose.\n"
        "  Configuració centralitzada de la DB.\n"
        "• Gestió: Migracions amb Alembic.\n"
        "  Actualitza l’esquema sense pèrdua de dades."
    )
    for p in content.paragraphs:
        p.font.size = Pt(22)
        p.font.color.rgb = text_color
    code = (
        "class User(UserMixin, db.Model):\n"
        "    id = db.Column(db.Integer, primary_key=True)\n"
        "    username = db.Column(db.String(80), unique=True)\n"
        "    password = db.Column(db.String(120))\n"
        "    theme = db.Column(db.String(20), default='light')"
    )
    add_code_box(slide, code, Inches(9), Inches(2), Inches(6), Inches(2), 14)
    add_image_with_border(slide, 'screenshots/render_dashboard_db.png', Inches(9), Inches(4.5), Inches(5), Inches(2), "Tauler de Control de la Base de Dades a Render")
    add_text_box(slide, "Estat de la base de dades PostgreSQL a Render", Inches(9), Inches(6.5), Inches(5), Inches(0.5), 16, text_color)
    add_text_box(slide, "Nota: Animació 'Slide In' per al codi", Inches(0.5), Inches(8), Inches(5), Inches(0.5), 12, text_color)

    # Slide 8: Frontend Features (Updated with Explanations)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    title = slide.shapes.title
    title.text = "Funcions del Frontend"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True

    features = [
        ("Multilingüe", "Suport per a ES, EN, CA.", Inches(1.5)),
        ("Entrada per Veu", "Webkit Speech Recognition.", Inches(3)),
        ("Suggerències", "Prompts dinàmics amb JS.", Inches(4.5)),
        ("Temes", "Clar i fosc amb CSS.", Inches(6))
    ]

    for i, (name, desc, top) in enumerate(features):
        add_icon(slide, Inches(1), top, Inches(0.6), secondary_color)
        add_text_box(slide, name, Inches(2), top, Inches(4), Inches(0.5), 26, accent_color, bold=True)
        explanation = {
            "Multilingüe": "Permet canviar l’idioma de la interfície.",
            "Entrada per Veu": "Converteix veu en text per a xats.",
            "Suggerències": "Ofereix prompts predefinits per a usuaris.",
            "Temes": "Personalitza l’aspecte visual de l’app."
        }[name]
        add_text_box(slide, f"{desc}\n  {explanation}", Inches(2), top + Inches(0.5), Inches(8), Inches(0.7), 20, text_color)
    code = (
        "<button id='voice-btn' class='btn btn-gradient'>\n"
        "    <i class='bi bi-mic-fill'></i>\n"
        "</button>"
    )
    add_code_box(slide, code, Inches(10), Inches(2), Inches(5), Inches(2), 16)
    add_text_box(slide, "Nota: Animació 'Fade In' per a funcions", Inches(0.5), Inches(8), Inches(5), Inches(0.5), 12, text_color)

    # Slide 9: Backend Features (Updated with Explanations)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    title = slide.shapes.title
    title.text = "Funcions del Backend"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True

    features = [
        ("Autenticació", "Flask-Login i bcrypt.", Inches(1.5)),
        ("Xat", "Integració amb OpenAI API.", Inches(3)),
        ("Historial", "Emmagatzematge a PostgreSQL.", Inches(4.5)),
        ("Configuració", "Personalització d’usuaris.", Inches(6))
    ]

    for i, (name, desc, top) in enumerate(features):
        add_icon(slide, Inches(1), top, Inches(0.6), secondary_color)
        add_text_box(slide, name, Inches(2), top, Inches(4), Inches(0.5), 26, accent_color, bold=True)
        explanation = {
            "Autenticació": "Protegix l’accés amb contrasenyes segures.",
            "Xat": "Genera respostes intel·ligents amb IA.",
            "Historial": "Guarda converses per a referència futura.",
            "Configuració": "Permet ajustar paràmetres de l’IA."
        }[name]
        add_text_box(slide, f"{desc}\n  {explanation}", Inches(2), top + Inches(0.5), Inches(8), Inches(0.7), 20, text_color)
    code = (
        "@app.route('/chat', methods=['POST'])\n"
        "def chat():\n"
        "    user_message = request.json.get('message')\n"
        "    response = openai.chat.completions.create(...)"
    )
    add_code_box(slide, code, Inches(10), Inches(2), Inches(5), Inches(2), 16)
    add_text_box(slide, "Nota: Animació 'Fade In' per a funcions", Inches(0.5), Inches(8), Inches(5), Inches(0.5), 12, text_color)

    # Slide 10: Challenges Faced (Updated with Explanations)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    title = slide.shapes.title
    title.text = "Dificultats Trobades"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True

    challenges = [
        ("Connexió DB", "Errors inicials de PostgreSQL.", Inches(2)),
        ("Variables", "Gestió de secrets a Render.", Inches(4)),
        ("Escalabilitat", "Optimització per a usuaris.", Inches(6))
    ]

    for i, (name, desc, top) in enumerate(challenges):
        add_icon(slide, Inches(1), top, Inches(0.6), secondary_color)
        add_text_box(slide, name, Inches(2), top, Inches(4), Inches(0.5), 26, accent_color, bold=True)
        explanation = {
            "Connexió DB": "Retards en establir connexió amb la DB.",
            "Variables": "Configuració segura de claus secretes.",
            "Escalabilitat": "Gestió de múltiples usuaris concurrents."
        }[name]
        add_text_box(slide, f"{desc}\n  {explanation}", Inches(2), top + Inches(0.5), Inches(8), Inches(0.7), 20, text_color)
    code = (
        "max_retries = 5\n"
        "while retry_count < max_retries:\n"
        "    try:\n"
        "        db.session.execute('SELECT 1')\n"
        "        db.create_all()\n"
        "        break\n"
        "    except OperationalError:\n"
        "        time.sleep(5)"
    )
    add_code_box(slide, code, Inches(9), Inches(2), Inches(6), Inches(4), 14)
    add_text_box(slide, "Nota: Animació 'Slide In' per a dificultats", Inches(0.5), Inches(8), Inches(5), Inches(0.5), 12, text_color)

    # Slide 11: Future Plans (Updated with Explanations)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    title = slide.shapes.title
    title.text = "Futur del Projecte"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True

    future = [
        ("IA Multimodal", "Suport per a imatges.", Inches(2)),
        ("Integracions", "Connexió amb Slack.", Inches(4)),
        ("Escalabilitat", "Milions d’usuaris.", Inches(6))
    ]

    for i, (name, desc, top) in enumerate(future):
        add_icon(slide, Inches(1), top, Inches(0.6), secondary_color)
        add_text_box(slide, name, Inches(2), top, Inches(4), Inches(0.5), 26, accent_color, bold=True)
        explanation = {
            "IA Multimodal": "Processament d’imatges i text combinats.",
            "Integracions": "Compatibilitat amb eines col·laboratives.",
            "Escalabilitat": "Suport per a una base d’usuaris massiva."
        }[name]
        add_text_box(slide, f"{desc}\n  {explanation}", Inches(2), top + Inches(0.5), Inches(8), Inches(0.7), 20, text_color)
    add_text_box(slide, "Nota: Animació 'Slide In' per a plans", Inches(0.5), Inches(8), Inches(5), Inches(0.5), 12, text_color)

    # Slide 12: Call to Action
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    set_slide_background(slide)

    title = slide.shapes.title
    title.text = "Uneix-te al Futur amb GarBotGPT!"
    title.text_frame.paragraphs[0].font.size = Pt(50)
    title.text_frame.paragraphs[0].font.color.rgb = accent_color
    title.text_frame.paragraphs[0].font.bold = True

    subtitle = slide.placeholders[1]
    subtitle.text = "Visita garbotgpt.com o contacta: support@garolacorp.com"
    subtitle.text_frame.paragraphs[0].font.size = Pt(28)
    subtitle.text_frame.paragraphs[0].font.color.rgb = text_color

    add_text_box(slide, "[Afegir codi QR per a garbotgpt.com]", Inches(6), Inches(4), Inches(4), Inches(0.5), 20, text_color, align=PP_ALIGN.CENTER)
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(6), Inches(4.5), Inches(2.5), Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = text_color
    add_text_box(slide, "Nota: Animació 'Pulse' per al codi QR", Inches(0.5), Inches(8), Inches(5), Inches(0.5), 12, text_color)

    # Save the presentation
    prs.save('GarBotGPT_Tech_Presentation_Updated_With_Screenshots.pptx')
    print("Presentació guardada com a 'GarBotGPT_Tech_Presentation_Updated_With_Screenshots.pptx'")

if __name__ == "__main__":
    create_garbotgpt_presentation()